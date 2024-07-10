import jwt
from datetime import datetime, timedelta
from functools import wraps
from flask import request, jsonify, current_app
from bcrypt import hashpw, gensalt
import uuid
import os
import math
from app.constants.roles import roles
import pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches


def generate_uuid():
    return str(uuid.uuid4())

def generate_defined_length_uuid(length=8):
    return str(uuid.uuid4()).replace('-', '')[:length]

def token_required(func):
    @wraps(func)
    def wrapper(*args, **kwargs):
        token = request.headers.get('Authorization')

        if not token:
            return jsonify({'message': 'Token is missing'}), 401

        try:
            secret_key = current_app.config['SECRET_KEY']
            data = jwt.decode(token, secret_key, algorithms=['HS256'])
            current_user = data
        except jwt.ExpiredSignatureError:
            return jsonify({'message': 'Token has expired'}), 401
        except jwt.InvalidTokenError:
            return jsonify({'message': 'Invalid token'}), 401

        return func(current_user=current_user, *args, **kwargs)

    return wrapper

def generate_jwt_token(email, role_id, user_id):
    payload = {
        'sub': email,
        'role_id': role_id,
        'id': user_id,
        'iat': datetime.utcnow(),
        'exp': datetime.utcnow() + timedelta(days=1)
    }

    secret_key = current_app.config['SECRET_KEY']

    token = jwt.encode(payload, secret_key, algorithm='HS256')
    return token

def superadmin_required(func):
    @wraps(func)
    def wrapper(*args, **kwargs):
        current_user = kwargs.get('current_user')

        if not current_user or current_user.get('role_id') != roles.get("SUPERADMIN"):
            return jsonify({'message': 'Access denied. Superadmin privileges required.'}), 403

        return func(*args, **kwargs)

    return wrapper

def admin_required(func):
    @wraps(func)
    def wrapper(*args, **kwargs):
        current_user = kwargs.get('current_user')

        if not current_user or (current_user.get('role_id') != roles.get("SUPERADMIN") and current_user.get('role_id') != roles.get("ADMIN")):
            return jsonify({'message': 'Access denied. Superadmin privileges required.'}), 403

        return func(*args, **kwargs)

    return wrapper

def clean_and_lower(value):
    if isinstance(value, str):
        return value.strip().lower()
    else:
        return value

def generate_bcrypt_hash(password):
    hashed_password = hashpw(password.encode('utf-8'), gensalt())
    return hashed_password.decode('utf-8')

def is_prod():
    return os.getenv('FLASK_ENV') == "production"

def haversine_distance(lat1, lon1, lat2, lon2):

    lat1, lon1, lat2, lon2 = map(math.radians, [lat1, lon1, lat2, lon2])

    dlat = lat2 - lat1
    dlon = lon2 - lon1
    a = math.sin(dlat / 2) ** 2 + math.cos(lat1) * math.cos(lat2) * math.sin(dlon / 2) ** 2
    c = 2 * math.asin(math.sqrt(a))
    r = 6371000  # Radius of Earth in meters. Use 6371 for kilometers.
    return c * r


def replace_image(prs, slide_index, image_index, new_image_path):
   
    slide = prs.slides[slide_index]
    
  
    image_count = 0
    old_image = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if image_count == image_index:
                old_image = shape
                break
            image_count += 1
    # Save position and size of the old image
    left = old_image.left
    top = old_image.top
    width = old_image.width
    height = old_image.height
    
    # Remove the old image
    slide.shapes._spTree.remove(old_image._element)
    
    # Add the new image in the same position
    slide.shapes.add_picture(new_image_path, left, top, width, height)
    
    
    
  