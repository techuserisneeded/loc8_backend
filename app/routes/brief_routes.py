from flask import Blueprint, request, jsonify, current_app, send_file
from datetime import datetime
from werkzeug.utils import secure_filename
import os
import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from io import BytesIO
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from app.utils.db_helper import query_db
from app.utils.helpers import token_required, generate_uuid, clean_and_lower, replace_image
from app.services.briefs import get_brief_details_by_brief_id, assign_brief_to_planners
from app.constants.roles import roles

brief_bp = Blueprint('brief', __name__)

ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg'}

osmo_logo_path = "./assets/Logo1.png"
no_image_path = "./assets/no_image.png"

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@brief_bp.route('/briefs', methods=['POST'])
@token_required
def createBrief(current_user):
    data = request.form

    current_user_id = current_user['id']

    if 'brand_logo' not in request.files:
        return jsonify({'message': 'Please provide brand logo'}), 400

    # assign data
    brief_id = generate_uuid();
    category = clean_and_lower(data['category'])
    brand_name = clean_and_lower(data['brand_name'])
    target_aud = clean_and_lower(data['target_aud'])
    camp_obj = clean_and_lower(data['camp_obj'])
    med_app = clean_and_lower(data['med_app'])
    is_immediate_camp = clean_and_lower(data['is_immediate_camp'])
    start_date = data.get("start_date")
    notes = data.get('notes')

    budgets = data.get("budgets")
    budgets = json.loads(budgets)

    if notes != None or notes != "":
        notes = clean_and_lower(notes) 

    if start_date != None or start_date != "":
        start_date = clean_and_lower(start_date) 

    file = request.files['brand_logo']

    # validate data
    if is_immediate_camp == 1 and not start_date:
        return jsonify({"message" : "Campaing is immdiate and we require start date!"}) , 400
    
    if file.filename == '':
        return jsonify({'message': 'No selected file'}), 400
    
    if start_date:
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        if start_date <= datetime.now():
            return jsonify({'message': 'Start date must be a future date'}), 400

    
    filename = brief_id + secure_filename(file.filename)
    file.save(os.path.join(current_app.config['UPLOAD_FOLDER'], filename))
    
    try:
        query_db("START TRANSACTION", ())

        brief_insert_q = """
                INSERT INTO `briefs`
                (
                    `brief_id`, `category`, `brand_name`, `brand_logo`, 
                    `target_audience`, `campaign_obj`, `media_approach`, `is_immediate_camp`, 
                    `start_date`, `notes`, `created_by_user_id`
                ) 
            VALUES (
                    %s, %s, %s, %s,
                    %s, %s, %s, %s,
                    %s, %s, %s
                )
        """

        query_db(brief_insert_q, (
            brief_id, category, brand_name, filename, 
            target_aud, camp_obj, med_app, is_immediate_camp,
            start_date, notes, current_user_id
        ))

        for budget in budgets:
            zone_id = budget.get("zone_id")
            state_id = budget.get("state_id")
            city_id = budget.get("city_id")
            budget_amt = budget.get("budget")
            budget_id = generate_uuid();

            budget_insert_q = """
                INSERT INTO `brief_budgets` (`budget_id`, `brief_id`, `zone_id`, `state_id`, `city_id`, `budget`)
                VALUES (%s, %s, %s, %s, %s, %s)
            """
            query_db(budget_insert_q, (budget_id, brief_id, zone_id, state_id, city_id, budget_amt))
        
        current_app.mysql.connection.commit()

        assign_brief_to_planners(brief_id)
        
        return jsonify({'message': 'Created successfully'}), 201
    except Exception as e:
        current_app.mysql.connection.rollback()
        print(e)
        return jsonify({'message': 'Something went wrong!'}), 500

@brief_bp.route('/briefs', methods=['GET'])
@token_required
def getBriefs(current_user):
    try:
        user_id = current_user['id']

        query = ""
        args = ()

        if current_user['role_id'] == roles.get("SUPERADMIN"):
            query = """
            SELECT b.*, bb.budget_id, bb.zone_id, bb.state_id, bb.city_id, zones.zone_name, states.state_name, cities.city_name, bb.budget,
            users.first_name, users.last_name
            FROM briefs b
            INNER JOIN brief_budgets bb ON b.brief_id = bb.brief_id
            INNER JOIN zones ON bb.zone_id = zones.zone_id
            INNER JOIN states ON bb.state_id = states.state_id
            INNER JOIN cities ON bb.city_id = cities.city_id
            INNER JOIN users ON b.created_by_user_id = users.id
        """
            
        else:

            query = """
                SELECT b.*, bb.budget_id, bb.zone_id, bb.state_id, bb.city_id, zones.zone_name, states.state_name, cities.city_name, bb.budget,
                users.first_name, users.last_name
                FROM briefs b
                INNER JOIN brief_budgets bb ON b.brief_id = bb.brief_id
                INNER JOIN zones ON bb.zone_id = zones.zone_id
                INNER JOIN states ON bb.state_id = states.state_id
                INNER JOIN cities ON bb.city_id = cities.city_id
                INNER JOIN users ON b.created_by_user_id = users.id
                WHERE b.created_by_user_id = %s
            """

            args = (user_id,)
            
        briefs_with_budgets = query_db(query, args)

        if briefs_with_budgets == None: 
            return jsonify([]), 200

        briefs_data = {}

        for row in briefs_with_budgets:
            brief_id = row['brief_id']
            if brief_id not in briefs_data:
                briefs_data[brief_id] = {
                    'brief_id': brief_id,
                    'category': row['category'],
                    'brand_name': row['brand_name'],
                    'brand_logo': row['brand_logo'],
                    'campaign_obj': row['campaign_obj'],
                    'start_date': row['start_date'],
                    'status': row['status'],
                    'first_name': row['first_name'],
                    'last_name': row['last_name'],
                    'budgets': []
                }
            if row['budget_id'] is not None:
                budget_data = {
                    'budget_id': row['budget_id'],
                    'zone_id': row['zone_id'],
                    'state_id': row['state_id'],
                    'city_id': row['city_id'],
                    'zone_name': row['zone_name'],
                    'state_name': row['state_name'],
                    'city_name': row['city_name'],
                    'budget': row['budget']
                }
                briefs_data[brief_id]['budgets'].append(budget_data)

        briefs_list = list(briefs_data.values())

        return jsonify(briefs_list), 200
    except Exception as e:
        print(e)
        return jsonify({'message': 'Something went wrong!'}), 500

@brief_bp.route('/briefs/<brief_id>', methods=['PUT'])
@token_required
def edit_brief(current_user, brief_id):
    data = request.form

    current_user_id = current_user['id']

    # assign data
    category = clean_and_lower(data['category'])
    brand_name = clean_and_lower(data['brand_name'])
    target_aud = clean_and_lower(data['target_aud'])
    camp_obj = clean_and_lower(data['camp_obj'])
    med_app = clean_and_lower(data['med_app'])
    is_immediate_camp = clean_and_lower(data['is_immediate_camp'])
    start_date = data.get("start_date")
    notes = data.get('notes')

    budgets = data.get("budgets")
    budgets = json.loads(budgets)

    map_img_filename = data.get("old_image", None)

    if notes != None or notes != "":
        notes = clean_and_lower(notes) 

    if start_date != None or start_date != "":
        start_date = clean_and_lower(start_date) 

    if 'brand_logo' in request.files:
        logo_image_file = request.files['brand_logo']
        if logo_image_file.filename != '':
            map_img_filename = generate_uuid() + secure_filename(logo_image_file.filename)
            logo_image_file.save(os.path.join(current_app.config['UPLOAD_FOLDER'], map_img_filename))

    # validate data
    if is_immediate_camp == 1 and not start_date:
        return jsonify({"message" : "Campaing is immdiate and we require start date!"}) , 400
    
    if start_date:
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        if start_date <= datetime.now():
            return jsonify({'message': 'Start date must be a future date'}), 400
    
    try:
        query_db("START TRANSACTION", ())

        brief_update_q = """
                UPDATE `briefs`
                SET
                    category=%s, brand_name=%s, 
                    brand_logo=%s, target_audience=%s, 
                    campaign_obj=%s, media_approach=%s, 
                    is_immediate_camp=%s, start_date=%s, 
                    notes=%s, created_by_user_id=%s
                WHERE
                    brief_id=%s
                
        """

        query_db(brief_update_q, (
            category, brand_name, 
            map_img_filename, target_aud, 
            camp_obj, med_app, 
            is_immediate_camp,start_date, 
            notes, current_user_id,
            brief_id
        ), commit=True)

        for budget in budgets:
            zone_id = budget.get("zone_id")
            state_id = budget.get("state_id")
            city_id = budget.get("city_id")
            budget_amt = budget.get("budget")

            budget_id = budget.get("budget_id")

            # only insert new budgets
            if not budget_id:
                budget_id = generate_uuid();

                budget_insert_q = """
                    INSERT INTO `brief_budgets` (`budget_id`, `brief_id`, `zone_id`, `state_id`, `city_id`, `budget`)
                    VALUES (%s, %s, %s, %s, %s, %s)
                """
                query_db(budget_insert_q, (budget_id, brief_id, zone_id, state_id, city_id, budget_amt))
        
        current_app.mysql.connection.commit()

        assign_brief_to_planners(brief_id)
        
        return jsonify({'message': 'Created successfully'}), 201
    except Exception as e:
        current_app.mysql.connection.rollback()
        print(e)
        return jsonify({'message': 'Something went wrong!'}), 500


@brief_bp.route('/briefs/<brief_id>', methods=['DELETE'])
@token_required
def deleteBrief(current_user, brief_id):

    brief_data = get_brief_details_by_brief_id(brief_id)

    # only admin and owner of the brief can delete 
    if(brief_data['created_by_user_id'] != current_user['id'] and current_user['role_id'] != roles.get("SUPERADMIN")):
        return jsonify({
            'message': "You cannot delete this brief!",
        }), 401
    
    try:
        query_db("START TRANSACTION")

        q = """
            DELETE FROM plans WHERE brief_id=%s
        """

        query_db(q, (brief_id,))
        
        q = """
            SELECT * FROM brief_budgets WHERE brief_id=%s
        """
        
        brief_budgets = query_db(q, (brief_id,))

        for bb in brief_budgets:
            budget_id = bb['budget_id']
            q = """
                DELETE FROM assigned_budgets WHERE budget_id=%s
            """
            query_db(q, (budget_id,))

        q = """
            DELETE FROM brief_budgets WHERE brief_id=%s
        """

        query_db(q, (brief_id,))

        q = """
            DELETE FROM briefs WHERE brief_id=%s
        """

        query_db(q, (brief_id,))

        current_app.mysql.connection.commit()

        brand_logo_path = os.path.join(current_app.config['UPLOAD_FOLDER'], brief_data['brand_logo'])

        if os.path.exists(brand_logo_path):
            os.remove(brand_logo_path)

        return jsonify({"message": "Deleted Successfully!"}), 200
    except Exception:
        query_db("ROLLBACK")
        return jsonify({"message": "Something went wrong!"}), 500

@brief_bp.route('/briefs/<brief_id>', methods=['GET'])
@token_required
def editBrief(current_user, brief_id):

    brief_data = get_brief_details_by_brief_id(brief_id)

    # # only admin and owner of the brief can delete 
    # if(brief_data['created_by_user_id'] != current_user['id'] and current_user['role_id'] != roles.get("SUPERADMIN")):
    #     return jsonify({
    #         'message': "You cannot edit this brief!",
    #     }), 401

    
    query = """
            SELECT bb.*, zones.zone_name, states.state_name, cities.city_name
            FROM brief_budgets bb
            INNER JOIN zones ON bb.zone_id = zones.zone_id
            INNER JOIN states ON bb.state_id = states.state_id
            INNER JOIN cities ON bb.city_id = cities.city_id
            WHERE bb.brief_id = %s
        """
    
    budgets = query_db(query, (brief_id,))

    brief_data["budgets"] = budgets

    return jsonify(brief_data), 200


@brief_bp.route('/planner', methods=['GET'])
@token_required
def getPlannerBriefs(current_user):
    user_id = current_user['id']

    assigned_brief_q = """
        SELECT b.*, bb.budget_id, bb.zone_id, bb.state_id, bb.city_id, zones.zone_name, states.state_name, cities.city_name, bb.budget,
        users.first_name, users.last_name
        FROM assigned_budgets ab
        INNER JOIN brief_budgets bb ON ab.budget_id = bb.budget_id
        INNER JOIN briefs b ON b.brief_id=bb.brief_id
        INNER JOIN zones ON bb.zone_id = zones.zone_id
        INNER JOIN states ON bb.state_id = states.state_id
        INNER JOIN cities ON bb.city_id = cities.city_id
        INNER JOIN users ON b.created_by_user_id = users.id
        WHERE 
            ab.user_id=%s
    """
    assigned_brief_args = (user_id,)

    briefs_with_budgets = query_db(assigned_brief_q, assigned_brief_args)

    if briefs_with_budgets == None: 
        return jsonify([]), 200

    briefs_data = {}

    for row in briefs_with_budgets:
        brief_id = row['brief_id']
        if brief_id not in briefs_data:
            briefs_data[brief_id] = {
                'brief_id': brief_id,
                'category': row['category'],
                'brand_name': row['brand_name'],
                'brand_logo': row['brand_logo'],
                'campaign_obj': row['campaign_obj'],
                'start_date': row['start_date'],
                'first_name': row['first_name'],
                'last_name': row['last_name'],
                'status': row['status'],
                'budgets': []
            }
        if row['budget_id'] is not None:
            budget_data = {
                'budget_id': row['budget_id'],
                'zone_id': row['zone_id'],
                'state_id': row['state_id'],
                'city_id': row['city_id'],
                'zone_name': row['zone_name'],
                'state_name': row['state_name'],
                'city_name': row['city_name'],
                'budget': row['budget']
            }
            briefs_data[brief_id]['budgets'].append(budget_data)

    briefs_list = list(briefs_data.values())

    return jsonify(briefs_list), 200

@brief_bp.route('/budgets/<budget_id>', methods=['GET'])
@token_required
def getBriefBudgetDetailsByBudgetId(current_user, budget_id):

    current_user_id = current_user['id']
    query = request.args

    average_speed_min = float(query.get("average_speed_min", 0) or 0)
    average_speed_max = float(query.get("average_speed_max", 99999) or 99999)
    
    query = """
            SELECT bb.budget_id, briefs.brief_id, bb.zone_id, bb.state_id, bb.city_id, zones.zone_name, states.state_name, cities.city_name, bb.budget, ab.status, ab.plan_ss 
            FROM assigned_budgets ab
            INNER JOIN brief_budgets bb ON ab.budget_id = bb.budget_id
            INNER JOIN briefs ON briefs.brief_id=bb.brief_id
            INNER JOIN zones ON bb.zone_id = zones.zone_id
            INNER JOIN states ON bb.state_id = states.state_id
            INNER JOIN cities ON bb.city_id = cities.city_id
            WHERE 
                    ab.budget_id=%s
                AND
                    ab.user_id=%s

        """
    
    video_query = """
        SELECT * FROM videofiles v
        WHERE v.zone_id=%s AND v.state_id=%s AND v.city_id=%s
        AND average_speed>=%s AND average_speed<=%s
    """
    
    plans_query = """
        SELECT plans.plan_id, plans.sr_no, b.* FROM plans
        INNER JOIN billboards b ON b.id=plans.billboard_id 
        WHERE budget_id=%s AND user_id=%s
        ORDER BY sr_no ASC
    """

    budget = query_db(query, (budget_id, current_user_id), True)

    if budget == None:
        return jsonify({}), 200

    plans = query_db(plans_query, (budget['budget_id'], current_user_id))
    videos = query_db(video_query, (budget['zone_id'], budget['state_id'], budget['city_id'], average_speed_min, average_speed_max))

    video_data = []

    if videos:
        for video in videos:
            q = """
                SELECT * FROM video_coordinates
                WHERE video_id=%s
            """

            args = (video['video_id'],)

            video_coords = query_db(q, args)

            if video_coords:
                video['coordinates'] = video_coords
            else:
                video['coordinates'] = []

            video_data.append(video)

    return jsonify({'budget':  budget, 'videos': video_data, 'plans': plans }), 200

@brief_bp.route('/briefs/<brief_id>/planner', methods=['GET'])
@token_required
def getBriefDetailsForPlanner(current_user, brief_id):

    brief_data = get_brief_details_by_brief_id(brief_id)
    user_id = current_user['id']
    
    query = """
            SELECT bb.budget_id, bb.zone_id, bb.state_id, bb.city_id, zones.zone_name, states.state_name, cities.city_name, bb.budget FROM assigned_budgets ab
            INNER JOIN brief_budgets bb ON ab.budget_id = bb.budget_id
            INNER JOIN zones ON bb.zone_id = zones.zone_id
            INNER JOIN states ON bb.state_id = states.state_id
            INNER JOIN cities ON bb.city_id = cities.city_id
            WHERE bb.brief_id=%s
                AND
                ab.user_id=%s
        """
    
    budgets = query_db(query, (brief_id, user_id))

    brief_data["budgets"] = budgets

    return jsonify(brief_data), 200

@brief_bp.route('/budgets/<budget_id>/<brief_id>/finish', methods=['PUT'])
@token_required
def finish_budget_plan(current_user, budget_id, brief_id):
    
    user_id = current_user['id']

    query = """
        UPDATE assigned_budgets
            SET status=2
        WHERE 
                budget_id=%s
            AND
                user_id=%s
    """

    query_db(query, (budget_id, user_id), True, True)

    query =  """
        SELECT count(budget_id) 
        FROM brief_budgets
        WHERE
            brief_id=%s
    """


    response = query_db(query, (brief_id,), True)

    # no of budget in the brief
    n_budgets_in_brief = response['count(budget_id)']

    query =  """
        SELECT count(assigned_budgets.budget_id)
        FROM assigned_budgets
        INNER JOIN brief_budgets
            ON brief_budgets.budget_id=assigned_budgets.budget_id
		WHERE
        		brief_budgets.brief_id=%s
            AND
            	assigned_budgets.status=2;
    """

    response = query_db(query, (brief_id,), True)

    # no of budget completed
    n_budgets_finished = response['count(assigned_budgets.budget_id)']

    if(n_budgets_in_brief == n_budgets_finished):
        query = """
            UPDATE briefs
                SET status=1
            WHERE
                brief_id=%s
        """

        query_db(query, (brief_id,), True, True)


    return jsonify({'message': "Plan updated!"})

@brief_bp.route("/briefs/<brief_id>/plans")
@token_required
def get_plans_by_brief_id(current_user, brief_id):

    q = """
            SELECT b.*, z.zone_name, s.state_name, c.city_name FROM plans p
            INNER JOIN billboards b
            ON b.id=p.billboard_id
            INNER JOIN brief_budgets bb 
            ON p.budget_id=bb.budget_id
            INNER JOIN zones z
            ON z.zone_id=bb.zone_id
            INNER JOIN states s
            ON s.state_id=bb.state_id
            INNER JOIN cities c
            ON c.city_id=bb.city_id
            WHERE p.brief_id=%s
        """
    args=(brief_id,)

    plan_details = query_db(q, args)
    

    if plan_details == None:
       return jsonify([]), 200
   
    return jsonify(plan_details), 200


@brief_bp.route('/assigned-budget/<budget_id>/image', methods=['PUT'])
@token_required
def update_assigned_budget_image(current_user, budget_id):

    current_user_id = current_user['id']

    map_img_file = request.files['map_img']
    map_img_filename = generate_uuid() + secure_filename(map_img_file.filename)

    map_img_file.save(os.path.join(current_app.config['UPLOAD_FOLDER'], map_img_filename))
    
    q = """
        UPDATE assigned_budgets
        SET plan_ss=%s
        WHERE user_id=%s AND budget_id=%s
    """

    query_db(q, (map_img_filename, current_user_id, budget_id), commit=True)

    return jsonify({'message':'image updated successfully!'}), 200



@brief_bp.route('/briefs/<brief_id>/download', methods=['GET'])
@token_required
def download_plan(current_user, brief_id):
    
    user_id = current_user['id']

    brief_details = get_brief_details_by_brief_id(brief_id)

    query = """
        SELECT bb.*, zones.zone_name, states.state_name, cities.city_name, ab.plan_ss
        FROM brief_budgets bb
        INNER JOIN zones ON bb.zone_id = zones.zone_id
        INNER JOIN states ON bb.state_id = states.state_id
        INNER JOIN cities ON bb.city_id = cities.city_id
        INNER JOIN assigned_budgets ab ON bb.budget_id=ab.budget_id
        WHERE bb.brief_id = %s
    """

    budgets = query_db(query, (brief_id,))
    
    presentation_path = os.path.join(current_app.config['ASSETS_FOLDER'], "new_input.pptx")
    brand_logo_path = os.path.join(current_app.config['UPLOAD_FOLDER'], brief_details['brand_logo'])

    if not os.path.exists(brand_logo_path):
        brand_logo_path = os.path.join(current_app.config['ASSETS_FOLDER'], "no_image.png")
        
    prs = Presentation(presentation_path)
    
    replace_image(prs, 0 , 0, brand_logo_path)
    
    for budget in budgets:

        # ALTER TABLE assigned_budgets ADD COLUMN plan_ss VARCHAR(300) NULL

        texts = [budget['zone_name'] + " | " + budget['state_name'] + " | " + budget['city_name']]

        if not budget['plan_ss']:
            budget_img_path = os.path.join(current_app.config['ASSETS_FOLDER'], "no_image.png")
        else:
            budget_img_path = os.path.join(current_app.config['UPLOAD_FOLDER'], budget['plan_ss'])

        secondPage(prs, texts, budget_img_path)
        
        query = """
            SELECT b.* FROM plans
            INNER JOIN billboards b ON b.id=plans.billboard_id
            WHERE plans.budget_id=%s
        """
        plans = query_db(query, (budget['budget_id'],))

        if not plans:
            continue

        for plan in plans:
            location  = plan['location']
            size_text = "Size: {}×{}".format(plan['height'], plan['width'])
            
            text = [location, size_text]
            
            site_image_path = current_app.config['UPLOAD_FOLDER'] +"/"+ plan['site_image']
            map_image_path = current_app.config['UPLOAD_FOLDER'] +"/"+ plan['map_image']

            if not os.path.exists(site_image_path):
                site_image_path = no_image_path

            if not os.path.exists(map_image_path):
                map_image_path = no_image_path

            print("image: ", site_image_path, file=sys.stdout)
            
            thirdPage(prs, text, site_image_path, map_image_path)
    
    #4th Page
    slide_layout = prs.slide_layouts[3]  
    slide = prs.slides.add_slide(slide_layout)
    
    presentation_bytes = BytesIO()
    prs.save(presentation_bytes)
    presentation_bytes.seek(0)

    return send_file(
        presentation_bytes,
        as_attachment=True,
        download_name="presentation.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    

#     prs = Presentation()

#     # slide dimensions for widescreen (16:9)
#     prs.slide_width = Inches(13.33)
#     prs.slide_height = Inches(7.5)

#     # slide with blank layout
#     blank_slide_layout = prs.slide_layouts[6]
#     intro_slide = prs.slides.add_slide(blank_slide_layout)

#     intro_slide = create_intro_slide(intro_slide, brand_logo_path, )

#     for budget in budgets:
#         area_slide = prs.slides.add_slide(blank_slide_layout)

#         area_slide = create_area_slide(area_slide, brand_logo_path, budget['zone_name'], budget['state_name'], budget['city_name'])

#         query = """
#             SELECT b.* FROM plans
#             INNER JOIN billboards b ON b.id=plans.billboard_id
#             WHERE plans.budget_id=%s
#         """

#         assets = query_db(query, (budget['budget_id'],))

#         if not assets:
#             continue

#         for plan in assets:
#             location  = plan['location']
#             size_text = "Size: {}×{}".format(plan['height'], plan['width'])

#             plan_slide = prs.slides.add_slide(blank_slide_layout)

#             # heading
#             title_shape = plan_slide.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(4.8), Inches(0.8))
#             title_text_frame = title_shape.text_frame
#             title_text_frame.text = "{}   {}".format(location, size_text)
#             title_text_frame.word_wrap = True
#             p = title_text_frame.paragraphs[0]
#             p.font.size = Inches(0.23)
#             p.font.color.rgb = RGBColor(255, 165, 0)

#             title_shape.fill.solid()
#             title_shape.fill.fore_color.rgb = RGBColor(0, 0, 255)

#             # brand logo
#             image_path = brand_logo_path
#             left_inch = Inches(8.5)
#             top_inch = Inches(0.2)
#             width_inch = Inches(2.5)
#             height_inch = Inches(1)
#             plan_slide.shapes.add_picture(image_path, left_inch, top_inch, width_inch, height_inch)

#             # Add osmo logo
#             image_path = osmo_logo_path
#             left_inch = Inches(11.3)
#             top_inch = Inches(0.1)
#             width_inch = Inches(2)
#             height_inch = Inches(1)
#             plan_slide.shapes.add_picture(image_path, left_inch, top_inch, width_inch, height_inch)

#             site_image_path = current_app.config['UPLOAD_FOLDER'] +"/"+ plan['site_image']
#             map_image_path = current_app.config['UPLOAD_FOLDER'] +"/"+ plan['map_image']

#             if not os.path.exists(site_image_path):
#                 site_image_path = no_image_path

#             if not os.path.exists(map_image_path):
#                 map_image_path = no_image_path

#             print("image: ", site_image_path, file=sys.stdout)

#             # site image
#             left_inch = Inches(1)
#             top_inch = Inches(1.2)
#             width_inch = Inches(11)
#             height_inch = Inches(6.2)
#             plan_slide.shapes.add_picture(site_image_path, left_inch, top_inch, width_inch, height_inch)

#             # Map Image
#             left_inch = Inches(9)
#             top_inch = Inches(5)
#             width_inch = Inches(4.2)
#             height_inch = Inches(2.5)
#             plan_slide.shapes.add_picture(map_image_path, left_inch, top_inch, width_inch, height_inch)




#     # creating thank you slide-
#     thank_u_slide = prs.slides.add_slide(blank_slide_layout)

#     # add logo
#     image_path = osmo_logo_path
#     left_inch = Inches(5.5)
#     top_inch = Inches(2.5)
#     width_inch = Inches(2.5)
#     height_inch = Inches(1.5)
#     thank_u_slide.shapes.add_picture(image_path, left_inch, top_inch, width_inch, height_inch)

#     # add title
#     title_shape = thank_u_slide.shapes.add_textbox(Inches(4.1), Inches(4), Inches(4.5), Inches(1))
#     title_text_frame = title_shape.text_frame
#     title_text_frame.text = "Let's create something irresistible"
#     p = title_text_frame.paragraphs[0]
#     p.font.size = Inches(0.4)
#     p.font.color.rgb = RGBColor(128, 128, 128)

#     presentation_bytes = BytesIO()
#     prs.save(presentation_bytes)
#     presentation_bytes.seek(0)


#     return send_file(
#         presentation_bytes,
#         as_attachment=True,
#         download_name="presentation.pptx",
#         mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
#     )

def secondPage(prs, text2ndPage, planner_image):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)

    for i, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        if i < len(text2ndPage):
            text_frame = shape.text_frame
            p = text_frame.paragraphs[0] 
            p.text = text2ndPage[i] + p.text  
            p.font.size = Pt(25)

            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p.alignment = PP_ALIGN.CENTER
            
    shape = slide.shapes[1]
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        old_image = shape
        sp = shape._element
        sp.getparent().remove(sp)
        left = old_image.left
        top = old_image.top
        width = old_image.width
        height = old_image.height
        slide.shapes.add_picture(planner_image, left, top, width, height)
        
        
def thirdPage(prs, texts, asset_image, map_image):
    slide_layout = prs.slide_layouts[2]  
    slide = prs.slides.add_slide(slide_layout)

    for i, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        if i < len(texts):
            text_frame = shape.text_frame
            p = text_frame.paragraphs[0] 
            p.text = texts[i] + p.text  

            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p.alignment = PP_ALIGN.CENTER
            
    shape = slide.shapes[2]
    # Check if the shape is a text box
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        old_image = shape
        sp = shape._element
        sp.getparent().remove(sp)
        left = old_image.left
        top = old_image.top
        width = old_image.width
        height = old_image.height
        slide.shapes.add_picture(asset_image, left, top, width, height)
    # Map Image
    left_inch = Inches(10.13)
    top_inch = Inches(5.34)
    width_inch = Inches(3)
    height_inch = Inches(2)
    slide.shapes.add_picture(map_image, left_inch, top_inch, width_inch, height_inch)
    
    return slide



# def create_intro_slide(intro_slide, brand_logo_path):
    
#     # brand logo
#     image_path = brand_logo_path
#     left_inch = Inches(5.5)
#     top_inch = Inches(1.3)
#     width_inch = Inches(2.5)
#     height_inch = Inches(1.5)
#     intro_slide.shapes.add_picture(image_path, left_inch, top_inch, width_inch, height_inch)

#     # heading
#     title_shape = intro_slide.shapes.add_textbox(Inches(4.5), Inches(3), Inches(4.5), Inches(1))
#     title_text_frame = title_shape.text_frame
#     title_text_frame.text = "On Street OOH"
#     p = title_text_frame.paragraphs[0]
#     p.font.size = Inches(0.7)
#     p.font.color.rgb = RGBColor(0, 0, 255)

#     title_shape.fill.solid()
#     title_shape.fill.fore_color.rgb = RGBColor(255, 165, 0)

#     title_shape = intro_slide.shapes.add_textbox(Inches(3.8), Inches(4.2), Inches(5.5), Inches(1))
#     title_text_frame = title_shape.text_frame
#     title_text_frame.text = "Recommendation"
#     p = title_text_frame.paragraphs[0]
#     p.font.size = Inches(0.7)
#     p.font.color.rgb = RGBColor(0, 0, 255)

#     title_shape.fill.solid()
#     title_shape.fill.fore_color.rgb = RGBColor(255, 165, 0)

#     # Add osmo logo
#     image_path = osmo_logo_path
#     left_inch = Inches(5.7)
#     top_inch = Inches(5.2)
#     width_inch = Inches(2)
#     height_inch = Inches(1.3)
#     intro_slide.shapes.add_picture(image_path, left_inch, top_inch, width_inch, height_inch)

#     return intro_slide

# def create_area_slide(intro_slide, brand_logo_path, zone="", state="", city=""):
    
#     # brand logo
#     image_path = brand_logo_path
#     left_inch = Inches(5.5)
#     top_inch = Inches(1.3)
#     width_inch = Inches(2.5)
#     height_inch = Inches(1.5)
#     intro_slide.shapes.add_picture(image_path, left_inch, top_inch, width_inch, height_inch)

#     areas = [zone, state, city]
#     top = 3

#     for area in areas:
#         # heading
#         title_shape = intro_slide.shapes.add_textbox(Inches(4.5), Inches(top), Inches(4.5), Inches(0.8))
#         title_text_frame = title_shape.text_frame
#         title_text_frame.text = area
#         p = title_text_frame.paragraphs[0]
#         p.font.size = Inches(0.4)
#         p.font.color.rgb = RGBColor(0, 0, 255)

#         title_shape.fill.solid()
#         title_shape.fill.fore_color.rgb = RGBColor(255, 165, 0)

#         top = top + 0.8

#     # Add osmo logo
#     image_path = osmo_logo_path
#     left_inch = Inches(5.7)
#     top_inch = Inches(5.2)
#     width_inch = Inches(2)
#     height_inch = Inches(1.3)
#     intro_slide.shapes.add_picture(image_path, left_inch, top_inch, width_inch, height_inch)

#     return intro_slide
